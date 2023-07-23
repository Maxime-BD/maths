import os
import json
from flask import Flask, render_template, request, redirect, url_for, flash, send_file
from werkzeug.utils import secure_filename
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

app = Flask(__name__)
app.secret_key = 'maberdeb'  # Replace with your secret key for Flask sessions

UPLOAD_FOLDER = 'upload'
ALLOWED_EXTENSIONS = {'txt', 'png', 'jpg', 'jpeg'}

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER


def allowed_file(filename):
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS


def add_title_slide(prs, title):
    # Choose the title slide layout
    slide_layout = prs.slide_layouts[5]

    # Add a slide with the title of the file without the extension
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title


def text_to_json(txt_file):
    # Read the lines from the text file and create a dictionary with slide names and calculations
    with open(txt_file, 'r') as file:
        lines = file.readlines()

    json_data = {}
    for idx, line in enumerate(lines):
        line = line.strip()

        if line:
            slide_name = f"Slide {idx + 1}"
            json_data[slide_name] = line

    return json_data


def create_ppt_from_json(json_filename, background_image, text_color):
    # Load the JSON data from the file
    with open(json_filename, 'r') as json_file:
        json_data = json.load(json_file)

    # Get the file name without the extension for the title slide
    file_title = os.path.splitext(os.path.basename(json_filename))[0]

    # Create a PowerPoint presentation with 16:9 aspect ratio
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # Add the first slide with the file name without the extension
    add_title_slide(prs, file_title)

    # Load the background image
    bg_img = background_image

    # Choose the color for the text (white or black)
    if text_color.lower() == 'white':
        font_color = RGBColor(255, 255, 255)  # White color
    else:
        font_color = RGBColor(0, 0, 0)  # Black color

    # Slide width and height
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Traverse the calculations in the JSON and create a slide for each
    for slide_name, calculation in json_data.items():
        slide_layout = prs.slide_layouts[6]  # Title and Content layout

        # Add a slide with the title "Slide X" and the calculation "1 + 5" in the content
        slide = prs.slides.add_slide(slide_layout)

        # Check if the slide has a title shape before setting its text
        if slide.shapes.title:
            slide.shapes.title.text = slide_name

        # Add the background image to the slide
        left = Inches(0)
        top = Inches(0)
        pic = slide.shapes.add_picture(bg_img, left, top, width=slide_width, height=slide_height)

        # Create a text box and set its properties to center the text
        textbox_left = Inches(2)
        textbox_top = (slide_height - Inches(4)) / 2  # Center vertically
        textbox_width = slide_width - Inches(4)
        textbox_height = slide_height - Inches(4)
        textbox = slide.shapes.add_textbox(textbox_left, textbox_top, textbox_width, textbox_height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True

        # Set the font properties for better quality and bigger text
        p = text_frame.add_paragraph()
        p.text = calculation
        p.font.size = Pt(48)  # Increase the font size as needed
        p.font.color.rgb = font_color
        p.alignment = PP_ALIGN.CENTER

    # Add a black slide at the end
    black_slide_layout = prs.slide_layouts[5]
    black_slide = prs.slides.add_slide(black_slide_layout)
    black_bg = black_slide.background
    black_fill = black_bg.fill
    black_fill.solid()
    black_fill.fore_color.rgb = RGBColor(0, 0, 0)  # Set the background to black

    # Save the presentation to a .pptx file with the same name as the JSON file
    pptx_filename = json_filename.replace('.json', '.pptx')

    # Delete the existing .pptx file if it already exists
    if os.path.exists(pptx_filename):
        os.remove(pptx_filename)

    prs.save(pptx_filename)
    return pptx_filename


@app.route('/download/<filename>', methods=['GET'])
def download_file(filename):
    pptx_file = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    return send_file(pptx_file, as_attachment=True)


@app.route('/', methods=['GET', 'POST'])
def index():
    if request.method == 'POST':
        # Check if the post request has the file part
        if 'file' not in request.files or 'image' not in request.files:
            flash('Please upload both a text file and an image.', 'error')
            return redirect(request.url)

        text_file = request.files['file']
        image_file = request.files['image']
        text_color = request.form['color']

        if text_file.filename == '' or image_file.filename == '':
            flash('Please select both a text file and an image.', 'error')
            return redirect(request.url)

        if text_file and allowed_file(text_file.filename) and image_file and allowed_file(image_file.filename):
            filename = secure_filename(text_file.filename)
            text_file.save(os.path.join(app.config['UPLOAD_FOLDER'], filename))
            image_filename = secure_filename(image_file.filename)
            image_file.save(os.path.join(app.config['UPLOAD_FOLDER'], image_filename))

            # Convert the text file to JSON
            json_data = text_to_json(os.path.join(app.config['UPLOAD_FOLDER'], filename))

            # Save the JSON data to a JSON file
            json_filename = os.path.join(app.config['UPLOAD_FOLDER'], f"{os.path.splitext(filename)[0]}.json")
            with open(json_filename, 'w') as json_file:
                json.dump(json_data, json_file)

            # Create the PowerPoint presentation from the JSON data
            ppt_filename = create_ppt_from_json(json_filename, os.path.join(app.config['UPLOAD_FOLDER'], image_filename),
                                                text_color)
            flash('Presentation created successfully!', 'success')

            return redirect(url_for('download_file', filename=os.path.basename(ppt_filename)))

    return render_template('cool.html')


if __name__ == "__main__":
    app.run(debug=True)
